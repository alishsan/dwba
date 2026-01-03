from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()

slides_content = [
    ("Development of LaBr₃(Ce)-Based Gamma Detector Modules for the NUMEN Project",
     "Uzbekistan–Italy Collaboration\n\nPresenter: [Your Name]\nAffiliation: New Uzbekistan University\nCollaborator: INFN-LNS (Italy)"),
    
    ("Scientific Background",
     "• Neutrinoless double-beta decay (0νββ): key to physics beyond the Standard Model\n"
     "• Nuclear matrix elements (NMEs) are crucial for 0νββ interpretation\n"
     "• The NUMEN project at INFN-LNS studies NMEs via heavy-ion reactions"),
    
    ("Motivation",
     "• NUMEN requires precise γ-ray detection for accurate NME extraction\n"
     "• Current systems limited in energy and timing resolution\n"
     "• Need for high-performance LaBr₃(Ce)-based detector modules"),
    
    ("Project Objectives",
     "• Develop and test 5 LaBr₃(Ce) detectors with PMTs\n"
     "• Integrate detectors into NUMEN γ-calorimeter\n"
     "• Enhance NUMEN’s sensitivity and precision\n"
     "• Build local expertise and infrastructure in Uzbekistan"),
    
    ("Detector Design & Concept",
     "• LaBr₃(Ce) properties:\n"
     "   – High light yield\n   – Fast decay time (~16 ns)\n   – Excellent resolution (~3% at 662 keV)\n"
     "• PMT coupling and signal readout design\n"
     "• Expected performance and calibration plans"),
    
    ("Assembly and Testing in Uzbekistan",
     "• Local laboratory setup and expertise\n"
     "• Assembly steps: crystal mounting, optical coupling, PMT connection\n"
     "• Testing:\n   – Energy calibration (¹³⁷Cs, ⁶⁰Co)\n   – Timing resolution and efficiency\n"
     "• Validation before shipment to INFN-LNS"),
    
    ("Integration at INFN-LNS",
     "• Installation in NUMEN γ-calorimeter array\n"
     "• Data acquisition compatibility\n"
     "• Contribution to precision NME measurements"),
    
    ("Applications in Uzbekistan",
     "• Radiation monitoring and environmental safety\n"
     "• Medical imaging (PET/SPECT)\n"
     "• Geological and mineral exploration\n"
     "• Border radiation control and security"),
    
    ("Educational and Industrial Impact",
     "• Student and researcher training at NewUU\n"
     "• Bilateral exchange: Uzbekistan ↔ Italy\n"
     "• Industry involvement in detector assembly/testing\n"
     "• Strengthens national high-tech sector"),
    
    ("Future Prospects & Commercialization",
     "• Project is applied, not directly commercial\n"
     "• Lays foundation for future technology transfer\n"
     "• Potential commercialization of detector systems\n"
     "• Enables access to international funding via NUMEN"),
    
    ("Strategic Outcomes",
     "• Integration of New Uzbekistan University into NUMEN consortium\n"
     "• Long-term detector technology development\n"
     "• Strengthened international scientific cooperation\n"
     "• Contribution to national innovation economy"),
    
    ("Summary & Outlook",
     "• 5 LaBr₃(Ce) detectors will enhance NUMEN’s precision\n"
     "• Builds local expertise and research capacity\n"
     "• Educational, industrial, and societal benefits\n"
     "• Strengthens Uzbekistan’s role in global nuclear research")
]

# Add slides
for title, content in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Save file
pptx_path = "/mnt/data/LaBr3_NUMEN_Project_Presentation.pptx"
prs.save(pptx_path)
pptx_path
